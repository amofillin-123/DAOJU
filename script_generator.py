#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import re
import sys
import glob
import logging
from io import BytesIO
from pptx import Presentation
from PIL import Image as PILImage
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak

# 配置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

class ScriptGenerator:
    def __init__(self):
        """初始化脚本生成器"""
        self.font_name = "STHeiti"
        logging.info("初始化 ScriptGenerator 完成")
        self.script_data = {
            "产品信息": {
                "名称": "",
                "链接": "",
                "主图": None
            },
            "产品卖点": [],
            "参考风格": [],
            "布景": {
                "布景风格": set(),  # 存储从"布景："后面提取的风格
                "拍摄场景": set(),  # 存储表格第一行的场景
                "拍摄场景": set()   # 将"场景内容"改为"拍摄场景"
            },
            "道具": {
                "装饰挂件": set(),
                "装饰材料": set(),
                "绿植类": set(),
                "辅助工具": set(),
                "套装类": set(),
                "场景布置": set()
            }
        }
        self.current_page = 0
        self.total_pages = 0
        self.current_file = ""

    def identify_slide_type(self, slide):
        """识别幻灯片类型"""
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text.strip())
        text = "\n".join(texts)
        
        if "产品信息" in text or any(keyword in text for keyword in ["产品链接", "产品名称"]):
            logging.info("识别为产品信息")
            return "产品信息页面"
        elif "产品卖点" in text:
            logging.info("识别为产品卖点页面")
            return "产品卖点页面"
        elif "参考风格" in text:
            logging.info("识别为参考风格页")
            return "参考风格页面"
        elif "拍摄思路" in text:
            logging.info("识别为拍摄思路页面")
            return "拍摄思路页面"
        
        logging.info("未识别页面类型")
        return None

    def process_product_info(self, slide):
        """处理产品信息页面"""
        logging.info("开始处理产品信息页面")
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                texts.append(text)
                
                # 提取产品名称
                if not self.script_data["产品信息"]["名称"] and not any(keyword in text.lower() for keyword in ["链接", "bgm", "产品信"]):
                    self.script_data["产品信息"]["名称"] = text
                    logging.info(f"提取产品名称: {text}")
                
                # 提取产品链接
                if "链接" in text.lower() and "http" in text:
                    link = re.search(r'https?://[^\s]+', text).group()
                    self.script_data["产品信息"]["链接"] = link
                    logging.info(f"提取产品链接: {link}")
            
            # 提取产品图片
            if hasattr(shape, "image") and self.script_data["产品信息"]["主图"] is None:
                try:
                    self.script_data["产品信息"]["主图"] = shape.image.blob
                    logging.info("成功提取产品图片")
                except Exception as e:
                    logging.error(f"提取产品图片失败: {e}")

    def process_selling_points(self, slide):
        """处理产品卖点页面"""
        logging.info("开始处理产品卖点页面")
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text and not text.startswith("0") and "产品卖点" not in text and "请输入" not in text:
                    texts.append(text)
                    logging.info(f"提取产品卖点: {text}")
        
        # 分割文本并添加到卖点列表
        for text in texts:
            points = re.split(r'[;；。]', text)
            for point in points:
                point = point.strip()
                if point and point not in self.script_data["产品卖点"]:
                    self.script_data["产品卖点"].append(point)

    def process_reference_style(self, slide):
        """处理参考风格页面"""
        logging.info("开始理参考风格页")
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    self.script_data["参考风格"].append(shape.image.blob)
                    logging.info("成功提取参考风格图片")
                except Exception as e:
                    logging.error(f"提取参考风格图片失败: {e}")

    def extract_scene_and_props(self, text, table_cells=None):
        """从文本中提取布景和道具信息"""
        logging.info(f"\n开始分析文本: {text}")
        
        # 提取布景风格
        if "布景：" in text:
            style_text = text.split("布景：")[1]
            if "道具：" in style_text:
                style_text = style_text.split("道具：")[0]
            style = style_text.strip()
            if style:
                self.script_data["布景"]["布景风格"].add(style)
                logging.info(f"提取布景风格: {style}")
        
        # 提取场景（如果提供了表格单元格）
        if table_cells and len(table_cells) > 0:
            # 检查第一个单元格否为"场景"（表头），如果是则跳过
            scene = table_cells[0].strip()
            if scene and scene != "场景":
                self.script_data["布景"]["拍摄场景"].add(scene)
                logging.info(f"提取拍摄场景: {scene}")
        
        # 提取道具信息
        if "道具：" in text:
            props_text = text.split("道具：")[1].strip()
            # 处理多种分隔符
            props = []
            for prop in re.split(r'[,，、\n]', props_text):
                # 处理空格分隔的道具
                props.extend([p.strip() for p in re.split(r'\s+', prop) if p.strip()])
            
            for prop in props:
                if prop:
                    # 处理带括号的说明文本
                    if '（' in prop:
                        prop = prop.split('（')[0].strip()
                    # 处理包含"等"、"及"、"或"的情况
                    for split_word in ["等", "及", "或"]:
                        if split_word in prop:
                            prop = prop.split(split_word)[0].strip()
                            break
                    if prop:
                        self.classify_prop(prop)

    def _classify_scene(self, scene, scene_keywords):
        """对场景进行分类"""
        if not scene or "拍摄思路" in scene:
            return
            
        # 检查每个类别的关键词
        for category, keywords in scene_keywords.items():
            if any(keyword in scene for keyword in keywords):
                # 对于拍摄角度，需要确保是真正的拍摄相关描述
                if category == "拍摄角度" and not any(k in scene for k in ["拍摄", "视角", "镜头", "特写", "远景", "近景"]):
                    continue
                    
                # 清理和规范化场景描述
                cleaned_scene = scene
                # 移除不必要的词语
                for word in ["在", "的", "地", "拍摄场景：场景"]:
                    cleaned_scene = cleaned_scene.replace(word, "")
                
                # 检查是否已经存在相同或相似的场景描述
                if not any(existing_scene in cleaned_scene or cleaned_scene in existing_scene 
                          for existing_scene in self.script_data["布景"][category]):
                    self.script_data["布景"][category].add(cleaned_scene)
                    logging.info(f"提取{category}: {cleaned_scene}")
                return  # 一个场景只分到一个类别
        
        # 如果没有匹配到任何类别，但包含场景相关词汇，归类到实景场景
        if any(word in scene for word in ["前", "后", "处", "边", "旁", "位"]):
            # 检查是否已经存在相同或相似的场景描述
            if not any(existing_scene in scene or scene in existing_scene 
                      for existing_scene in self.script_data["布景"]["实景场景"]):
                self.script_data["布景"]["实景场景"].add(scene)
                logging.info(f"提取实景场景: {scene}")

    def classify_prop(self, prop):
        """对道具进行分类"""
        # 道具类别关键词映射
        prop_categories = {
            "装饰挂件": ["挂", "吊", "装饰", "饰品", "球", "花环"],
            "装饰材料": ["纸", "布", "条", "带", "绳", "丝带", "藤条"],
            "绿植类": ["树", "枝", "叶", "花", "草", "绿植"],
            "辅助工具": ["剪刀", "胶", "钉", "针", "工具"],
            "套装类": ["套装", "套件", "组合"],
            "场景布置": ["桌布", "背景", "道具", "布景"]
        }
        
        # 检查每个类别
        for category, keywords in prop_categories.items():
            if any(keyword in prop for keyword in keywords):
                self.script_data["道具"][category].add(prop)
                logging.info(f"提取{category}道具: {prop}")
                return  # 个道具只分到一个类别
        
        # 如果没有匹配到何类别，归类到场景布置
        self.script_data["道具"]["场景布置"].add(prop)
        logging.info(f"提取场景布置道具: {prop}")

    def process_shooting_idea(self, slide):
        """处理拍摄思路页面"""
        logging.info("开始处理拍摄思路页面")
        
        # 处理表格
        for shape in slide.shapes:
            if shape.has_table:
                self.process_table(shape.table)
            elif hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    self.extract_scene_and_props(text)

    def process_table(self, table):
        """处理表格内容"""
        if not table.rows:
            return
        
        # 获取第一行的所有单元格内容
        first_row = table.rows[0]
        if not first_row.cells:
            return
        
        # 跳过第一个单元格(场景标题)，处理后面的单元格
        for i in range(1, len(first_row.cells)):
            scene = first_row.cells[i].text.strip()
            if scene:  # 确保不是空单元格
                self.script_data["布景"]["拍摄场景"].add(scene)
                logging.info(f"提取拍摄场景: {scene}")

    def process_file(self, filename):
        """处理单个PPTX文件"""
        try:
            prs = Presentation(filename)
            self.total_pages = len(prs.slides)
            print(f"\n总页数: {self.total_pages}\n")
            
            # 初始化数据结构
            self.script_data = {
                "产品信息": {
                    "名称": "",
                    "链接": "",
                    "主图": None
                },
                "产品卖点": [],
                "参考风格": [],
                "布景": {
                    "布景风格": set(),
                    "拍摄场景": set()
                },
                "道具": {
                    "装饰挂件": set(),
                    "装饰材料": set(),
                    "绿植类": set(),
                    "辅助工具": set(),
                    "套装类": set(),
                    "场景布置": set()
                }
            }

            # 处理每一页
            for i, slide in enumerate(prs.slides, 1):
                self.current_page = i
                print(f"{'-'*30}")
                print(f"处理第 {i} 页:")
                
                # 识别页面类型并处理
                slide_type = self.identify_slide_type(slide)
                if slide_type:
                    print(f"识别为: {slide_type}")
                    
                    if slide_type == "产品信息页面":
                        self.process_product_info(slide)
                    elif slide_type == "产品卖点页面":
                        self.process_selling_points(slide)
                    elif slide_type == "参考风格页面":
                        self.process_reference_style(slide)
                    elif slide_type == "拍摄思路页面":
                        self.process_shooting_idea(slide)
                else:
                    print("未识别页面类型")
            
            # 生成输出文件名
            output_filename = os.path.splitext(filename)[0] + "_拍摄需求.pdf"
            self.generate_pdf(output_filename)
            
            # 打印提取内容摘要
            self.print_summary()
            
            print(f"\n{'='*50}")
            print(f"完成处理: {filename}")
            print(f"生成的文件：{output_filename}")
            print(f"{'='*50}\n")
            
        except Exception as e:
            logging.error(f"处理文件 {filename} 时发生错误: {e}")
            raise

    def generate_pdf(self, output_filename):
        """生成PDF文档"""
        doc = SimpleDocTemplate(
            output_filename,
            pagesize=A4,
            rightMargin=30,
            leftMargin=30,
            topMargin=30,
            bottomMargin=30
        )
        
        # 创建样式
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontName=self.font_name,
            fontSize=28,
            textColor=colors.black,  # 主标题使用黑色
            spaceAfter=30,
            alignment=1  # 居中
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontName=self.font_name,
            fontSize=20,
            textColor=colors.HexColor('#0066CC'),  # 其他标题使用蓝色
            spaceAfter=20
        )
        
        # 添加子标题样式
        subheading_style = ParagraphStyle(
            'CustomSubHeading',
            parent=styles['Heading3'],
            fontName=self.font_name,
            fontSize=16,
            textColor=colors.black,
            spaceAfter=10,
            spaceBefore=10
        )
        
        normal_style = ParagraphStyle(
            'CustomNormal',
            parent=styles['Normal'],
            fontName=self.font_name,
            fontSize=12,
            leading=20,
            spaceAfter=10
        )
        
        bullet_style = ParagraphStyle(
            'CustomBullet',
            parent=styles['Normal'],
            fontName=self.font_name,
            fontSize=12,
            leading=20,
            leftIndent=20,
            spaceAfter=5
        )
        
        # 构建每页的内容
        page_contents = []
        
        # 第1页：产品信息
        page1 = []
        page1.append(Paragraph("拍摄需求文档", title_style))
        page1.append(Spacer(1, 20))
        page1.append(Paragraph("1. 产品信息", heading_style))
        page1.append(Paragraph(f"产品名称：{self.script_data['产品信息']['名称']}", normal_style))
        page1.append(Paragraph(f"产品链接：{self.script_data['产品信息']['链接']}", normal_style))
        
        if self.script_data["产品信息"]["主图"]:
            try:
                img_data = BytesIO(self.script_data["产品信息"]["主图"])
                img = PILImage.open(img_data)
                img = img.convert('RGB')
                
                # 计算缩放后的尺寸
                aspect = img.height / img.width
                new_width = 4 * inch  # 设置图片宽度为4英寸
                new_height = new_width * aspect
                
                # 添加产品图片
                page1.append(Spacer(1, 10))
                page1.append(Image(img_data, width=new_width, height=new_height))
            except Exception as e:
                logging.error(f"处理产品图片时发生错误: {e}")
        
        page_contents.append(page1)
        
        # 第2页：产品卖点
        page2 = []
        page2.append(Paragraph("2. 产品卖点", heading_style))
        for point in self.script_data["产品卖点"]:
            # 处理数字加粗
            highlighted_point = self.highlight_numbers(point)
            page2.append(Paragraph(f"• {highlighted_point}", bullet_style))
        page_contents.append(page2)
        
        # 第3页：参考风格
        page3 = []
        page3.append(Paragraph("3. 参考风格", heading_style))
        if self.script_data["参考风格"]:
            # 创建三列布局的表格
            images_data = []
            current_row = []
            column_width = doc.width / 3
            
            for i, img_blob in enumerate(self.script_data["参考风格"]):
                try:
                    img_data = BytesIO(img_blob)
                    img = PILImage.open(img_data)
                    img = img.convert('RGB')
                    
                    # 计算缩放后的尺寸
                    aspect = img.height / img.width
                    new_width = column_width
                    new_height = new_width * aspect
                    
                    current_row.append(Image(img_data, width=new_width, height=new_height))
                    
                    if len(current_row) == 3:
                        images_data.append(current_row)
                        current_row = []
                except Exception as e:
                    logging.error(f"处理参考风格图片时发生错误: {e}")
            
            # 处理最后一行不足三列的情况
            if current_row:
                while len(current_row) < 3:
                    current_row.append("")
                images_data.append(current_row)
            
            if images_data:
                # 创建表格并设置样式
                table = Table(images_data, colWidths=[column_width] * 3)
                table.setStyle(TableStyle([
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
                    ('LEFTPADDING', (0, 0), (-1, -1), 5),
                    ('RIGHTPADDING', (0, 0), (-1, -1), 5),
                    ('TOPPADDING', (0, 0), (-1, -1), 5),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
                ]))
                page3.append(table)
        page_contents.append(page3)
        
        # 第4页：布景要求
        page4 = []
        page4.append(Paragraph("4. 布景要求", heading_style))
        for category, scenes in self.script_data["布景"].items():
            if scenes:  # 只显示有内容的分类
                page4.append(Paragraph(category, subheading_style))
                for scene in sorted(scenes):
                    # 为数字添加红色加粗样式
                    highlighted_scene = self.highlight_numbers(scene)
                    page4.append(Paragraph(f"• {highlighted_scene}", bullet_style))
                page4.append(Spacer(1, 10))
        page_contents.append(page4)
        
        # 第5页：道具清单
        page5 = []
        page5.append(Paragraph("5. 道具清单", heading_style))
        
        # 使用计数器来生成序号
        counter = 1
        for category, props in self.script_data["道具"].items():
            if props:  # 只显示有内容的分类
                page5.append(Paragraph(category, subheading_style))
                sorted_props = sorted(props)  # 对道具进行排序
                for prop in sorted_props:
                    # 为数字添加红色加粗样式
                    highlighted_prop = self.highlight_numbers(prop)
                    # 使用数字序号替代原来的圆点
                    page5.append(Paragraph(f"{counter}. {highlighted_prop}", bullet_style))
                    counter += 1
                page5.append(Spacer(1, 10))
        
        page_contents.append(page5)
        
        # 构建最终文档内容
        story = []
        for i, page_content in enumerate(page_contents):
            # 添加页面内容
            story.extend(page_content)
            # 在每页之后添加分页符（最后一页除外）
            if i < len(page_contents) - 1:
                story.append(PageBreak())
        
        # 生成PDF
        doc.build(story)

    def print_summary(self):
        """打印提取内容摘要"""
        print("\n" + "="*30)
        print("提取内容摘要:")
        print("\n1. 产品信息:")
        print(f"   名称: {self.script_data['产品信息']['名称']}")
        print(f"   链接: {self.script_data['产品信息']['链接']}")
        print("   图片: " + ("已提取" if self.script_data['产品信息']['主图'] else "未提取"))
        
        print("\n2. 产品卖点:")
        for point in self.script_data['产品卖点']:
            print(f"   - {point}")
            
        print(f"\n3. 参考风格片: {len(self.script_data['参考风格'])} 张")
        
        print("\n4. 布景信息:")
        for category in ["布景风格", "拍摄场景"]:
            if self.script_data["布景"][category]:
                print(f"   {category}:")
                for item in self.script_data["布景"][category]:
                    print(f"   - {item}")
                    
        print("\n5. 道具信息:")
        for category in ["装饰挂件", "装饰材料", "绿植类", "辅助工具", "套装类", "场景布置"]:
            if self.script_data["道具"][category]:
                print(f"   {category}:")
                for prop in self.script_data["道具"][category]:
                    print(f"   - {prop}")

    def highlight_numbers(self, text):
        """为文本中的数字添加红色加粗样式"""
        if not text:
            return text
        
        # 定义中文数字映射
        chinese_nums = {
            '零': 0, '一': 1, '二': 2, '三': 3, '四': 4, '五': 5, '六': 6, '七': 7, '八': 8, '九': 9,
            '十': 10, '百': 100, '千': 1000, '万': 10000, '亿': 100000000
        }
        
        # 创建中文数字的正则表达式模式
        chinese_pattern = f"[{''.join(chinese_nums.keys())}]+"
        
        # 先处理阿拉伯数字
        text = re.sub(r'(\d+)', r'<b><font color="red">\1</font></b>', text)
        
        # 再处理中文数字
        def replace_chinese_num(match):
            return f'<b><font color="red">{match.group(0)}</font></b>'
        
        text = re.sub(chinese_pattern, replace_chinese_num, text)
        
        return text

if __name__ == "__main__":
    # 注册字体
    font_path = "/System/Library/Fonts/STHeiti Light.ttc"
    try:
        pdfmetrics.registerFont(TTFont("STHeiti", font_path))
        print(f"\n成功加载字体: STHeiti")
    except Exception as e:
        logging.error(f"加载字体失败: {e}")
        sys.exit(1)
    
    # 查找当前目录下的所PPTX文件
    pptx_files = glob.glob("*.pptx")
    
    if not pptx_files:
        print("错误：当前目录下没有找到PPTX文件")
        exit(1)
    
    print("\n找到以下PPTX文:")
    for file in pptx_files:
        print(f"- {file}")
    
    # 创建脚本生成器实例
    generator = ScriptGenerator()
    
    # 处理每个PPTX文件
    for filename in pptx_files:
        print(f"\n{'='*50}")
        print(f"开始处理文件: {filename}")
        print(f"{'='*50}")
        generator.process_file(filename) 