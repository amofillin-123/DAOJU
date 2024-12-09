# -*- coding: utf-8 -*-
import os
import re
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.lib.enums import TA_CENTER
from PIL import Image as PILImage
from io import BytesIO
import glob

class ScriptGenerator:
    def __init__(self):
        """初始化脚本生成器"""
        self.font_name = "STHeiti"
        self.script_data = {
            "产品信息": {"名称": "", "链接": "", "图片": None},
            "产品卖点": [],
            "参考风格": [],
            "布景": {
                "主要场景": set(),
                "辅助场景": set(),
                "氛围布置": set(),
                "拍摄角度": set()
            },
            "道具": {
                "装饰挂件类": set(),
                "装饰材料类": set(),
                "绿植类": set(),
                "辅助工具类": set(),
                "套装类": set(),
                "场景布置相关": set()
            }
        }
        
    def process_slide(self, slide):
        """处理单个幻灯片内容"""
        is_selling_point_slide = False
        is_product_info_slide = False
        is_reference_style_slide = False
        
        # 检查页面类型
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if "02| 产品卖点" in text:
                    is_selling_point_slide = True
                    break
                elif "产品信息" in text:
                    is_product_info_slide = True
                    break
                elif "参考风格" in text:
                    is_reference_style_slide = True
                    break
        
        # 如果是参考风格页面，提取所有图片
        if is_reference_style_slide:
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        self.script_data["参考风格"].append(shape.image.blob)
                    except Exception as e:
                        print(f"警告：提取参考风格图片时出错: {str(e)}")
        
        # 如果是产品卖点页面，提取所有实际的卖点内容
        if is_selling_point_slide:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text and not text.startswith("02|") and "产品卖点" not in text:
                        self.script_data["产品卖点"].append(text)
        
        # 处理其他内容
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if not text:
                    continue
                
                # 提取产品名称（从任何页面）
                if "J7-" in text or "迷你圣诞树装饰品套装" in text:
                    self.script_data["产品信息"]["名称"] = text.strip()
                
                # 只从产品信息页面提取链接
                if is_product_info_slide and "amazon.com" in text:
                    # 使用正则表达式提取URL
                    url_match = re.search(r'https?://[^\s\]]+', text)
                    if url_match:
                        self.script_data["产品信息"]["链接"] = url_match.group(0)
                
                # 提取布景和道具信息
                self.extract_props_and_scenes(text)
            
            # 处理图片
            if hasattr(shape, "image"):
                # 检查是否产品图片（这里假设第一张图片是产品图）
                if self.script_data["产品信息"]["图片"] is None:
                    try:
                        self.script_data["产品信息"]["图片"] = shape.image.blob
                    except Exception as e:
                        print(f"警告：提取图片时出错: {str(e)}")
            
            # 处理表格
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            self.extract_props_and_scenes(cell_text)

    def extract_props_and_scenes(self, text):
        """从文本中提取道具和布景信息"""
        if not text:
            return
            
        # 清理说明性文字
        text = re.sub(r'(道具[：:]\s*|画面[：:]\s*|布景[：:]\s*)', '', text)
        
        # 布景分类规则
        scene_categories = {
            "主要场景": [
                (r'怪诞风(?:圣诞)?桌面', lambda m: "怪诞风格圣诞桌面"),
                (r'壁炉(?:区域)?', lambda m: "壁炉区域"),
                (r'歪脖子圣诞树(?:区域)?', lambda m: "歪脖子圣诞树区域")
            ],
            "辅助场景": [
                (r'窗户', lambda m: "窗户"),
                (r'门把手', lambda m: "门把手"),
                (r'墙面', lambda m: "墙面")
            ],
            "氛围布置": [
                (r'怪诞风格', lambda m: "怪诞风格布置"),
                (r'圣诞氛', lambda m: "圣诞氛围布置"),
                (r'家中的任何部分', lambda m: "家居整体布置")
            ],
            "拍摄角度": [
                (r'过肩拍摄', lambda m: "过肩拍摄视角"),
                (r'多角度展示', lambda m: "多角度展示"),
                (r'镜头扫过', lambda m: "横向推移视角")
            ]
        }
        
        # 应用布景分类规则
        for category, patterns in scene_categories.items():
            for pattern, formatter in patterns:
                matches = re.finditer(pattern, text)
                for match in matches:
                    scene = formatter(match)
                    if scene:
                        self.script_data["布景"][category].add(scene)
        
        # 道具分类规则
        prop_categories = {
            "装饰挂件类": [
                (r'(\d+)\s*个([^，。\n]+?(?:圣诞球|挂牌|装饰品))', lambda m: f"{m.group(1)}个{m.group(2).strip()}")
            ],
            "装饰材料类": [
                (r'一卷([^，。\n]+?(?:丝带))', lambda m: f"一卷{m.group(1).strip()}")
            ],
            "绿植类": [
                (r'小型翠绿色歪脖子树', lambda m: m.group(0)),
                (r'圣诞树(?!装饰)', lambda m: "圣诞树"),
                (r'圣诞树枝', lambda m: "圣诞树枝")
            ],
            "辅助工具类": [
                (r'剪刀', lambda m: "剪刀"),
                (r'红色绿色桌纸', lambda m: "红色绿色桌纸"),
                (r'藤条', lambda m: "藤条")
            ],
            "套装类": [
                (r'完整的圣诞装饰套装', lambda m: "完整的圣诞装饰套装"),
                (r'迷你圣诞树装饰品套装', lambda m: "你圣诞饰品套装")
            ],
            "场景布置相关": [
                (r'壁炉', lambda m: "壁炉"),
                (r'花环', lambda m: "花环")
            ]
        }
        
        # 应用道具分类规则
        for category, patterns in prop_categories.items():
            for pattern, formatter in patterns:
                matches = re.finditer(pattern, text)
                for match in matches:
                    prop = formatter(match)
                    if prop:
                        self.script_data["道具"][category].add(prop)

    def process_selling_points(self, text):
        """处理产品卖点信息，将其分类整理"""
        selling_points = {
            "产品套装内容": [],
            "使用场景": [],
            "产品特点": []
        }
        
        # 遍历所有卖点
        for point in self.script_data["产品卖点"]:
            # 清理文本
            point = point.strip()
            if not point:
                continue
            
            # 移除可能的标题前缀
            point = re.sub(r'^.*?[:|：]\s*', '', point)
            
            # 分解复合句子
            parts = re.split(r'[;；]', point)
            for part in parts:
                part = part.strip()
                if not part:
                    continue
                
                # 处理产品套装内容
                if re.search(r'\d+\s*个|一卷', part):
                    # 分解套装内容为单独的项目
                    items = re.split(r'[,，、]', part)
                    for item in items:
                        item = item.strip()
                        if re.search(r'\d+\s*个|一卷', item):
                            selling_points["产品套装内容"].append(item)
                
                # 处理使用场景
                elif re.search(r'适合|可以|挂在|用于', part):
                    # 提取具体场景
                    scenes = re.findall(r'[圣诞树|树枝|窗户|门把手]+', part)
                    for scene in scenes:
                        if scene:
                            selling_points["使用场景"].append(f"适合挂在{scene}")
                
                # 处理产品特点
                elif re.search(r'设计|魅力|氛围|多样性', part):
                    selling_points["产品特点"].append(part)
        
        return selling_points

    def highlight_numbers(self, text):
        """为文本中的数字添加红色加粗样式"""
        # 匹配"数字+个"的模式
        pattern = r'(\d+\s*个)'
        parts = re.split(pattern, text)
        result = []
        for i, part in enumerate(parts):
            if i % 2 == 0:  # 非数字部分
                result.append(part)
            else:  # 数字部分
                result.append(f'<font color="red"><b>{part}</b></font>')
        return ''.join(result)

    def create_image_flow_layout(self, images, max_width=500, spacing=10):
        """创建瀑布流布局"""
        if not images:
            return None
            
        # 计算每行可以放置的图片数量（这里设置为3列）
        columns = 3
        column_width = (max_width - (columns - 1) * spacing) / columns
        
        # 处理所有图片
        processed_images = []
        for img_data in images:
            try:
                img = PILImage.open(BytesIO(img_data))
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # 计算缩放后的尺寸
                aspect = img.height / img.width
                new_width = column_width
                new_height = column_width * aspect
                
                # 限制最大高度
                max_height = 200  # 设置最大高度
                if new_height > max_height:
                    new_height = max_height
                    new_width = new_height / aspect
                
                img.thumbnail((new_width, new_height), PILImage.Resampling.LANCZOS)
                
                img_buffer = BytesIO()
                img.save(img_buffer, format='JPEG')
                img_buffer.seek(0)
                
                img_reader = Image(img_buffer)
                img_reader.drawWidth = new_width
                img_reader.drawHeight = new_height
                processed_images.append(img_reader)
            except Exception as e:
                print(f"警告：处理参考风格图片时出错: {str(e)}")
        
        # 创建表格数据
        table_data = []
        row = []
        for i, img in enumerate(processed_images):
            row.append(img)
            if len(row) == columns:
                table_data.append(row)
                row = []
        if row:  # 处理最后一行不满的情况
            while len(row) < columns:
                row.append('')
            table_data.append(row)
        
        # 创建表格样式
        style = TableStyle([
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('LEFTPADDING', (0, 0), (-1, -1), spacing/2),
            ('RIGHTPADDING', (0, 0), (-1, -1), spacing/2),
            ('TOPPADDING', (0, 0), (-1, -1), spacing/2),
            ('BOTTOMPADDING', (0, 0), (-1, -1), spacing/2),
        ])
        
        # 创建表格
        table = Table(table_data)
        table.setStyle(style)
        return table

    def generate_pdf(self, output_filename):
        """生成PDF文档"""
        doc = SimpleDocTemplate(
            output_filename,
            pagesize=A4,
            rightMargin=36,
            leftMargin=36,
            topMargin=36,
            bottomMargin=36
        )
        
        # 创建样式
        styles = getSampleStyleSheet()
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=28,
            spaceAfter=20,
            spaceBefore=20,
            textColor=colors.black,  # 主标题使用黑色
            fontName=self.font_name,
            leading=34,
            alignment=TA_CENTER
        )
        
        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=20,
            spaceAfter=16,
            spaceBefore=16,
            textColor=colors.HexColor('#0066CC'),  # 其他标题使用蓝色
            fontName=self.font_name,
            leading=24,
            alignment=TA_CENTER
        )
        
        subheading_style = ParagraphStyle(
            'CustomSubHeading',
            parent=styles['Heading3'],
            fontSize=14,
            spaceAfter=10,
            spaceBefore=10,
            textColor=colors.black,
            fontName=self.font_name,
            leading=18
        )
        
        content_style = ParagraphStyle(
            'CustomContent',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=8,
            leading=16,
            textColor=colors.black,
            fontName=self.font_name,
            firstLineIndent=20
        )
        
        bullet_style = ParagraphStyle(
            'CustomBullet',
            parent=styles['Normal'],
            fontSize=12,
            spaceAfter=8,
            leading=16,
            textColor=colors.black,
            fontName=self.font_name,
            leftIndent=20,
            firstLineIndent=0
        )
        
        # 构建文档内容
        story = []
        
        # 第1页：产品信息
        story.append(Paragraph("拍摄需求文档", title_style))  # 主标题使用黑色
        story.append(Spacer(1, 30))
        story.append(Paragraph('<font color="#0066CC">产品信息</font>', heading_style))  # 二级标题使用蓝色
        if any(self.script_data["产品信息"].values()):
            if self.script_data["产品信息"]["名称"]:
                story.append(Paragraph("产品名称：", subheading_style))
                story.append(Paragraph(self.script_data["产品信息"]["名称"], content_style))
                story.append(Spacer(1, 10))
            
            if self.script_data["产品信息"]["链接"]:
                story.append(Paragraph("产品链接：", subheading_style))
                story.append(Paragraph(self.script_data["产品信息"]["链接"], content_style))
                story.append(Spacer(1, 10))
            
            if self.script_data["产品信息"]["图片"]:
                try:
                    img_data = BytesIO(self.script_data["产品信息"]["图片"])
                    img = PILImage.open(img_data)
                    
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    # 调整图片大小，使其适合页
                    max_size = (500, 500)
                    img.thumbnail(max_size, PILImage.Resampling.LANCZOS)
                    
                    img_buffer = BytesIO()
                    img.save(img_buffer, format='JPEG')
                    img_buffer.seek(0)
                    
                    img_reader = Image(img_buffer)
                    aspect = img.height / img.width
                    img_reader.drawWidth = 400
                    img_reader.drawHeight = 400 * aspect
                    
                    story.append(img_reader)
                except Exception as e:
                    print(f"警告：添加产品图片时出错: {str(e)}")
        
        story.append(PageBreak())
        
        # 第2页：产品卖点
        story.append(Paragraph('<font color="#0066CC">产品卖点</font>', heading_style))
        if self.script_data["产品卖点"]:
            selling_points = self.process_selling_points("\n".join(self.script_data["产品卖点"]))
            for category, points in selling_points.items():
                if points:
                    story.append(Paragraph(category, subheading_style))
                    for point in points:
                        highlighted_point = self.highlight_numbers(point)
                        story.append(Paragraph("• " + highlighted_point, bullet_style))
                    story.append(Spacer(1, 10))
        
        story.append(PageBreak())
        
        # 第3页：参考风格
        story.append(Paragraph('<font color="#0066CC">参考风格</font>', heading_style))
        if self.script_data["参考风格"]:
            story.append(Spacer(1, 20))
            image_flow = self.create_image_flow_layout(self.script_data["参考风格"])
            if image_flow:
                story.append(image_flow)
        
        story.append(PageBreak())
        
        # 第4页：布景要求
        story.append(Paragraph('<font color="#0066CC">布景要求</font>', heading_style))
        for category, scenes in self.script_data["布景"].items():
            if scenes:  # 只显示有内容的分类
                story.append(Paragraph(category, subheading_style))
                for scene in sorted(scenes):
                    # 为数字添加红色加粗样式
                    highlighted_scene = self.highlight_numbers(scene)
                    story.append(Paragraph("• " + highlighted_scene, bullet_style))
                story.append(Spacer(1, 10))
        
        story.append(PageBreak())
        
        # 第5页：道具清单
        story.append(Paragraph('<font color="#0066CC">道具清单</font>', heading_style))
        for category, props in self.script_data["道具"].items():
            if props:  # 只显示有内容的分类
                story.append(Paragraph(category, subheading_style))
                for prop in sorted(props):
                    # 为数字添加红色加粗样式
                    highlighted_prop = self.highlight_numbers(prop)
                    story.append(Paragraph("• " + highlighted_prop, bullet_style))
                story.append(Spacer(1, 10))
        
        # 生成PDF
        doc.build(story)

    def process_file(self, filename):
        """处理单个PPTX文件"""
        print(f"\n开始处理文件: {filename}")
        try:
            prs = Presentation(filename)
            # 处理所有幻灯片
            for slide in prs.slides:
                self.process_slide(slide)
            
            # 注册字体
            font_path = "STHeiti Light.ttc"
            pdfmetrics.registerFont(TTFont(self.font_name, font_path))
            print(f"成功加载字体: {self.font_name}")
            
            # 生成PDF
            self.generate_pdf("拍摄需求.pdf")
            print(f"完成处理: {filename}")
            print(f"生成的文件：拍摄需求.pdf")
        except Exception as e:
            print(f"处理文件 {filename} 时出错: {str(e)}")

    def generate_script(self):
        """生成脚本文档"""
        # 获取当前目录下的所有pptx文件
        pptx_files = glob.glob("*.pptx")
        
        if not pptx_files:
            print("错误：当前目录下没有找到PPTX文件")
            return
            
        for filename in pptx_files:
            self.process_file(filename)
            # 置数据，准备处理下一个文件
            self.script_data = {
                "产品信息": {"名称": "", "链接": "", "图片": None},
                "产品卖点": [],
                "参考风格": [],
                "布景": {
                    "主要场景": set(),
                    "辅助场景": set(),
                    "氛围布置": set(),
                    "拍摄角度": set()
                },
                "道具": {
                    "装饰挂件类": set(),
                    "装饰材料类": set(),
                    "绿植类": set(),
                    "辅助工具类": set(),
                    "套装类": set(),
                    "场景布置相关": set()
                }
            }

if __name__ == "__main__":
    generator = ScriptGenerator()
    generator.generate_script() 